from __future__ import annotations

import csv
from io import StringIO
from pathlib import Path

from amr_preprocess.models import (
    DocClass,
    ExtractedDocument,
    ExtractedTable,
    FigureAsset,
    RawDocument,
    TextBlock,
)


def extract_docx(raw: RawDocument) -> ExtractedDocument:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required for .docx") from exc

    doc = Document(raw.bytes_path)
    blocks: list[TextBlock] = []
    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name if para.style else "") or ""
        kind = "heading" if style.lower().startswith("heading") else "paragraph"
        blocks.append(TextBlock(block_id=f"p{i}", type=kind, text=text, page=1))

    tables: list[ExtractedTable] = []
    for ti, table in enumerate(doc.tables):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        headers, body = ([rows[0]], rows[1:]) if rows else ([], [])
        tables.append(
            ExtractedTable(
                table_id=f"t{ti}",
                page=1,
                headers=headers,
                rows=body,
                extraction_method="native",
                confidence=1.0,
            )
        )

    figures: list[FigureAsset] = []
    return ExtractedDocument(
        doc_id=raw.doc_id,
        parent_doc_id=raw.parent_doc_id,
        source_uri=raw.source_uri,
        filename=raw.filename,
        doc_class=DocClass.DOCX,
        mime_type=raw.mime_type,
        blocks=blocks,
        tables=tables,
        figures=figures,
        page_count=1,
        metadata={**raw.metadata, "paragraphs": len(doc.paragraphs)},
        warnings=[],
    )


def extract_pptx(raw: RawDocument) -> ExtractedDocument:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for .pptx") from exc

    prs = Presentation(raw.bytes_path)
    blocks: list[TextBlock] = []
    tables: list[ExtractedTable] = []
    figures: list[FigureAsset] = []
    warnings: list[str] = []

    for si, slide in enumerate(prs.slides, start=1):
        for shi, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text = "\n".join(
                    p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
                )
                if text:
                    kind = "heading" if shi == 0 else "paragraph"
                    blocks.append(
                        TextBlock(
                            block_id=f"s{si}_sh{shi}",
                            type=kind,
                            text=text,
                            page=si,
                        )
                    )
            if shape.has_table:
                tbl = shape.table
                rows = [
                    [cell.text.strip() for cell in row.cells] for row in tbl.rows
                ]
                headers, body = ([rows[0]], rows[1:]) if rows else ([], [])
                tables.append(
                    ExtractedTable(
                        table_id=f"s{si}_t{shi}",
                        page=si,
                        headers=headers,
                        rows=body,
                        extraction_method="native",
                        confidence=1.0,
                    )
                )
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                figures.append(
                    FigureAsset(
                        figure_id=f"s{si}_img{shi}",
                        page=si,
                        kind="unknown",
                    )
                )

    return ExtractedDocument(
        doc_id=raw.doc_id,
        parent_doc_id=raw.parent_doc_id,
        source_uri=raw.source_uri,
        filename=raw.filename,
        doc_class=DocClass.PPTX,
        mime_type=raw.mime_type,
        blocks=blocks,
        tables=tables,
        figures=figures,
        page_count=len(prs.slides),
        metadata={**raw.metadata, "slides": len(prs.slides)},
        warnings=warnings,
    )


def extract_xlsx(raw: RawDocument) -> ExtractedDocument:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("openpyxl is required for .xlsx") from exc

    wb = load_workbook(raw.bytes_path, data_only=True, read_only=True)
    tables: list[ExtractedTable] = []
    blocks: list[TextBlock] = []
    for si, name in enumerate(wb.sheetnames):
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append(["" if c is None else str(c).strip() for c in row])
        rows = [r for r in rows if any(c for c in r)]
        blocks.append(
            TextBlock(block_id=f"sheet{si}", type="heading", text=name, page=si + 1)
        )
        headers, body = ([rows[0]], rows[1:]) if rows else ([], [])
        tables.append(
            ExtractedTable(
                table_id=f"sheet{si}",
                page=si + 1,
                caption=name,
                headers=headers,
                rows=body,
                extraction_method="native",
                confidence=1.0,
            )
        )
    wb.close()
    return ExtractedDocument(
        doc_id=raw.doc_id,
        parent_doc_id=raw.parent_doc_id,
        source_uri=raw.source_uri,
        filename=raw.filename,
        doc_class=DocClass.SHEET,
        mime_type=raw.mime_type,
        blocks=blocks,
        tables=tables,
        page_count=len(tables),
        metadata=raw.metadata,
    )


def extract_csv(raw: RawDocument) -> ExtractedDocument:
    text = Path(raw.bytes_path).read_text(encoding="utf-8", errors="replace")
    reader = csv.reader(StringIO(text))
    rows = [[c.strip() for c in row] for row in reader]
    headers, body = ([rows[0]], rows[1:]) if rows else ([], [])
    table = ExtractedTable(
        table_id="csv0",
        page=1,
        headers=headers,
        rows=body,
        extraction_method="native",
        confidence=1.0,
    )
    return ExtractedDocument(
        doc_id=raw.doc_id,
        parent_doc_id=raw.parent_doc_id,
        source_uri=raw.source_uri,
        filename=raw.filename,
        doc_class=DocClass.SHEET,
        mime_type=raw.mime_type,
        tables=[table],
        page_count=1,
        metadata=raw.metadata,
    )
