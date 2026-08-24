from __future__ import annotations

from email.message import EmailMessage
from pathlib import Path

import pymupdf
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt


def write_samples(dest: Path) -> list[Path]:
    dest = Path(dest)
    dest.mkdir(parents=True, exist_ok=True)
    written = [
        _report_pdf(dest / "acme-annual-report.pdf"),
        _deck_pdf(dest / "acme-earnings-deck.pdf"),
        _docx(dest / "account-notes.docx"),
        _pptx(dest / "qbr-slides.pptx"),
        _xlsx(dest / "pipeline.xlsx"),
        _csv(dest / "risks.csv"),
        _eml(dest / "review-pack.eml", dest / "risks.csv"),
        _text(dest / "call-notes.txt"),
    ]
    return written


def write_eval_fixtures(dest: Path) -> list[Path]:
    dest = Path(dest)
    dest.mkdir(parents=True, exist_ok=True)
    cases = []

    report = dest / "report-pdf"
    report.mkdir(exist_ok=True)
    _report_pdf(report / "acme-annual-report.pdf")
    (report / "expected.json").write_text(
        _json(
            {
                "filename": "acme-annual-report.pdf",
                "doc_class": "report",
                "must_contain": ["Acme Corp FY25 Annual Report", "subscription and services revenue"],
                "tables": [
                    {
                        "headers": [["Segment", "FY25", "FY24"]],
                        "rows": [["Subscription", "1200", "1000"], ["Services", "300", "250"]],
                    }
                ],
            }
        ),
        encoding="utf-8",
    )
    cases.append(report)

    deck = dest / "deck-pdf"
    deck.mkdir(exist_ok=True)
    _deck_pdf(deck / "acme-earnings-deck.pdf")
    (deck / "expected.json").write_text(
        _json(
            {
                "filename": "acme-earnings-deck.pdf",
                "doc_class": "deck",
                "must_contain": ["Q2 FY26 Results"],
                "tables": [
                    {
                        "headers": [["Metric", "Q2", "YoY"]],
                        "rows": [["Revenue", "412", "9%"], ["OpMargin", "31%", "120bps"]],
                    }
                ],
            }
        ),
        encoding="utf-8",
    )
    cases.append(deck)

    notes = dest / "docx-notes"
    notes.mkdir(exist_ok=True)
    _docx(notes / "account-notes.docx")
    (notes / "expected.json").write_text(
        _json(
            {
                "filename": "account-notes.docx",
                "doc_class": "docx",
                "must_contain": ["Churn risk", "Expansion"],
                "tables": [
                    {
                        "headers": [["Risk", "Severity"]],
                        "rows": [["Churn", "High"], ["Support backlog", "Medium"]],
                    }
                ],
            }
        ),
        encoding="utf-8",
    )
    cases.append(notes)

    sheet = dest / "xlsx-pipeline"
    sheet.mkdir(exist_ok=True)
    _xlsx(sheet / "pipeline.xlsx")
    (sheet / "expected.json").write_text(
        _json(
            {
                "filename": "pipeline.xlsx",
                "doc_class": "sheet",
                "must_contain": ["Opportunities"],
                "tables": [
                    {
                        "headers": [["Account", "Stage", "Amount"]],
                        "rows": [["Globex", "Commit", "250000"]],
                    }
                ],
            }
        ),
        encoding="utf-8",
    )
    cases.append(sheet)

    mail = dest / "eml-pack"
    mail.mkdir(exist_ok=True)
    csv_path = dest / "_tmp_risks.csv"
    _csv(csv_path)
    _eml(mail / "review-pack.eml", csv_path, attachment_name="risks.csv")
    csv_path.unlink(missing_ok=True)
    (mail / "expected.json").write_text(
        _json(
            {
                "filename": "review-pack.eml",
                "doc_class": "email",
                "must_contain": ["AMR review pack"],
                "tables": [],
            }
        ),
        encoding="utf-8",
    )
    cases.append(mail)
    return cases


def _json(data: dict) -> str:
    import json

    return json.dumps(data, indent=2)


def _report_pdf(path: Path) -> Path:
    doc = pymupdf.open()
    page = doc.new_page(width=612, height=792)
    page.insert_text((72, 72), "Acme Corp FY25 Annual Report", fontsize=18)
    page.insert_text(
        (72, 110),
        "This 10-K style report summarizes subscription and services revenue.",
        fontsize=11,
    )
    page.insert_text((72, 160), "Segment          FY25      FY24", fontsize=11)
    page.insert_text((72, 180), "Subscription     1200      1000", fontsize=11)
    page.insert_text((72, 200), "Services          300       250", fontsize=11)
    page.insert_text((72, 240), "See the financial statements for full detail.", fontsize=11)
    # make it look dense enough to classify as a report
    y = 280
    for i in range(18):
        page.insert_text(
            (72, y),
            f"Narrative paragraph {i}: remaining performance obligation and cash flow notes.",
            fontsize=10,
        )
        y += 16
    doc.set_metadata({"creator": "Adobe InDesign 19.0", "producer": "Adobe PDF Library"})
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))
    doc.close()
    return path


def _deck_pdf(path: Path) -> Path:
    doc = pymupdf.open()
    page = doc.new_page(width=792, height=612)
    page.insert_text((48, 48), "Q2 FY26 Results", fontsize=22)
    page.insert_text((48, 90), "Metric     Q2      YoY", fontsize=14)
    page.insert_text((48, 120), "Revenue    412      9%", fontsize=14)
    page.insert_text((48, 150), "OpMargin   31%      120bps", fontsize=14)
    page.insert_text((48, 200), "Guidance unchanged. See appendix.", fontsize=12)
    doc.set_metadata({"creator": "Microsoft PowerPoint", "producer": "PDFMaker 16 for PowerPoint"})
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))
    doc.close()
    return path


def _docx(path: Path) -> Path:
    doc = Document()
    doc.add_heading("Account review notes", level=1)
    doc.add_paragraph("Expansion opportunity in EMEA. Churn risk on the core seat renewal.")
    table = doc.add_table(rows=3, cols=2)
    table.rows[0].cells[0].text = "Risk"
    table.rows[0].cells[1].text = "Severity"
    table.rows[1].cells[0].text = "Churn"
    table.rows[1].cells[1].text = "High"
    table.rows[2].cells[0].text = "Support backlog"
    table.rows[2].cells[1].text = "Medium"
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))
    return path


def _pptx(path: Path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    tf = box.text_frame
    tf.text = "QBR snapshot"
    tf.paragraphs[0].font.size = Pt(28)
    rows, cols = 3, 2
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(6), Inches(2)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "NPS"
    table.cell(1, 1).text = "42"
    table.cell(2, 0).text = "Seats"
    table.cell(2, 1).text = "1800"
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


def _xlsx(path: Path) -> Path:
    wb = Workbook()
    ws = wb.active
    ws.title = "Opportunities"
    ws.append(["Account", "Stage", "Amount"])
    ws.append(["Globex", "Commit", 250000])
    ws.append(["Initech", "Propose", 80000])
    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(path))
    return path


def _csv(path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text("risk,owner\nChurn,AM\nSecurity review,SE\n", encoding="utf-8")
    return path


def _eml(path: Path, attachment: Path, attachment_name: str | None = None) -> Path:
    msg = EmailMessage()
    msg["From"] = "ae@example.com"
    msg["To"] = "am@example.com"
    msg["Subject"] = "AMR review pack"
    msg.set_content("Please process the attached risks.csv for this week's review.")
    msg.add_attachment(
        attachment.read_bytes(),
        maintype="text",
        subtype="csv",
        filename=attachment_name or attachment.name,
    )
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(msg.as_bytes())
    return path


def _text(path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        "Call notes\n\nCustomer asked about Agentforce rollout and seat expansion.\n",
        encoding="utf-8",
    )
    return path
